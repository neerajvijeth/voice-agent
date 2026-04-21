"""
voice_agent/rag/document_loader.py

Parse .docx and .pptx files into text chunks with source metadata.
Each chunk carries its source filename and location (slide number or paragraph range)
so the agent can tell the user exactly where the information comes from.

Chunking strategy:
  - Extract raw text per logical unit (slide, paragraph group, page)
  - Re-chunk with overlapping windows for better retrieval at chunk boundaries
  - chunk_size and chunk_overlap are configurable in config.py
"""

import os
import re
from dataclasses import dataclass
from config import RAG_CHUNK_SIZE, RAG_CHUNK_OVERLAP


@dataclass
class Chunk:
    """A text chunk with source metadata."""
    text: str
    source_file: str
    location: str       # e.g. "Slide 3" or "Paragraphs 5-10"


# ── Generic overlapping chunker ──────────────────────────────────────────────

def _split_into_sentences(text: str) -> list[str]:
    """Split text into sentences using simple regex."""
    # Split on sentence-ending punctuation followed by whitespace or end-of-string
    parts = re.split(r'(?<=[.!?])\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]


def chunk_text(text: str, source_file: str, location: str,
               chunk_size: int = RAG_CHUNK_SIZE,
               chunk_overlap: int = RAG_CHUNK_OVERLAP) -> list[Chunk]:
    """
    Split text into overlapping chunks on sentence boundaries.

    - Accumulates sentences until chunk_size is reached
    - Backs up by chunk_overlap characters for the next chunk
    - Ensures no information is lost at chunk boundaries
    """
    sentences = _split_into_sentences(text)
    if not sentences:
        return []

    chunks = []
    current_sentences = []
    current_len = 0

    for sentence in sentences:
        current_sentences.append(sentence)
        current_len += len(sentence) + 1  # +1 for space

        if current_len >= chunk_size:
            chunk_text_str = " ".join(current_sentences).strip()
            if chunk_text_str:
                chunks.append(Chunk(
                    text=chunk_text_str,
                    source_file=source_file,
                    location=location,
                ))

            # Keep overlap: walk backwards until we have ~chunk_overlap chars
            overlap_len = 0
            overlap_sentences = []
            for s in reversed(current_sentences):
                overlap_len += len(s) + 1
                overlap_sentences.insert(0, s)
                if overlap_len >= chunk_overlap:
                    break

            current_sentences = overlap_sentences
            current_len = sum(len(s) + 1 for s in current_sentences)

    # Flush remaining
    if current_sentences:
        chunk_text_str = " ".join(current_sentences).strip()
        if chunk_text_str:
            chunks.append(Chunk(
                text=chunk_text_str,
                source_file=source_file,
                location=location,
            ))

    return chunks


# ── File loaders ─────────────────────────────────────────────────────────────

def load_docx(path: str) -> list[Chunk]:
    """Extract text from a .docx file, chunked with overlap."""
    from docx import Document
    doc = Document(path)
    fname = os.path.basename(path)

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []

    full_text = " ".join(paragraphs)
    return chunk_text(full_text, source_file=fname, location="Document")


def load_pptx(path: str) -> list[Chunk]:
    """Extract text from a .pptx file, chunked per slide with overlap."""
    from pptx import Presentation
    prs = Presentation(path)
    fname = os.path.basename(path)
    all_chunks = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)

        if texts:
            slide_text = " ".join(texts)
            slide_chunks = chunk_text(
                slide_text,
                source_file=fname,
                location=f"Slide {slide_num}",
            )
            all_chunks.extend(slide_chunks)

    return all_chunks


def load_pdf(path: str) -> list[Chunk]:
    """Extract text from a PDF file, chunked with overlap."""
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    fname = os.path.basename(path)
    all_chunks = []

    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            page_chunks = chunk_text(
                text.strip(),
                source_file=fname,
                location=f"Page {page_num}",
            )
            all_chunks.extend(page_chunks)

    return all_chunks


def load_txt(path: str) -> list[Chunk]:
    """Load a plain text file, chunked with overlap."""
    fname = os.path.basename(path)
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read().strip()

    if not content:
        return []

    return chunk_text(content, source_file=fname, location="Document")


def load_directory(directory: str) -> list[Chunk]:
    """Load all supported documents from a directory."""
    all_chunks = []

    if not os.path.isdir(directory):
        print(f"[RAG] Documents directory not found: {directory}")
        return []

    for fname in sorted(os.listdir(directory)):
        fpath = os.path.join(directory, fname)
        if not os.path.isfile(fpath):
            continue

        try:
            if fname.endswith(".docx"):
                chunks = load_docx(fpath)
            elif fname.endswith(".pptx"):
                chunks = load_pptx(fpath)
            elif fname.endswith(".pdf"):
                chunks = load_pdf(fpath)
            elif fname.endswith(".txt"):
                chunks = load_txt(fpath)
            else:
                continue

            print(f"[RAG] Loaded {fname}: {len(chunks)} chunks")
            all_chunks.extend(chunks)
        except Exception as e:
            print(f"[RAG] Error loading {fname}: {e}")

    print(f"[RAG] Total: {len(all_chunks)} chunks from {directory}")
    return all_chunks
